#!/usr/bin/env python3
from __future__ import annotations

import argparse
import base64
import hashlib
import html
import json
import math
import re
import shutil
import subprocess
import tempfile
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

TAG_A = "{" + NS["a"] + "}"
TAG_P = "{" + NS["p"] + "}"
TAG_R = "{" + NS["r"] + "}"

CANVAS = {"width": 1600, "height": 900}
FONT_FALLBACK = "'TikTok Display', 'Inter', 'PingFang SC', 'Microsoft YaHei', Arial, sans-serif"


@dataclass(frozen=True)
class Transform:
    ox: float = 0.0
    oy: float = 0.0
    sx: float = 1.0
    sy: float = 1.0

    def apply_rect(self, x: float, y: float, w: float, h: float) -> tuple[float, float, float, float]:
        return self.ox + x * self.sx, self.oy + y * self.sy, w * self.sx, h * self.sy


class Converter:
    def __init__(self, pptx_path: Path, out_dir: Path):
        self.pptx_path = pptx_path
        self.out_dir = out_dir
        self.assets_dir = out_dir / "assets" / "pptx-media"
        self.prs = Presentation(str(pptx_path))
        self.scale_x = CANVAS["width"] / self.prs.slide_width
        self.scale_y = CANVAS["height"] / self.prs.slide_height
        self.theme = self._load_theme_colors()
        self.asset_by_hash: dict[str, dict] = {}
        self.assets: list[dict] = []

    def run(self) -> None:
        if self.out_dir.exists():
            shutil.rmtree(self.out_dir)
        self.assets_dir.mkdir(parents=True, exist_ok=True)

        slide_entries = []
        slide_html = []
        for index, slide in enumerate(self.prs.slides, 1):
            fragments = []
            title = self._slide_title(slide, index)
            element_count = 0

            for shape in self._background_shapes(slide):
                rendered = self._render_shape(shape, Transform(), source="layout")
                if rendered:
                    fragments.append(rendered)
                    element_count += rendered.count('class="ls-element')

            for shape in slide.shapes:
                rendered = self._render_shape(shape, Transform(), source="slide")
                if rendered:
                    fragments.append(rendered)
                    element_count += rendered.count('class="ls-element')

            if element_count == 0:
                fragments.append('<div class="ls-element ls-shape" style="left:0;top:0;width:1600px;height:900px;background:#fff;"></div>')
                element_count = 1

            slide_entries.append(
                {
                    "index": index,
                    "title": title,
                    "layout": slide.slide_layout.name,
                    "elementCount": element_count,
                }
            )
            slide_html.append(self._slide_section(index, title, "\n".join(fragments)))

        manifest = {
            "source": self.pptx_path.name,
            "slideCount": len(slide_html),
            "canvas": CANVAS,
            "slides": slide_entries,
            "assets": self.assets,
        }
        (self.out_dir / "source-deck-manifest.json").write_text(
            json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8"
        )
        (self.out_dir / "lark-design-guidelines.html").write_text(
            self._deck_html("\n".join(slide_html)), encoding="utf-8"
        )

    def _background_shapes(self, slide) -> Iterable:
        master = slide.slide_layout.slide_master
        for shape in master.shapes:
            if not self._is_placeholder(shape):
                yield shape
        for shape in slide.slide_layout.shapes:
            if not self._is_placeholder(shape):
                yield shape

    def _render_shape(self, shape, transform: Transform, source: str) -> str:
        if self._is_hidden(shape):
            return ""
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return self._render_group(shape, transform, source)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return self._render_picture(shape, transform, source)
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            return self._render_line(shape, transform, source)
        if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            return self._render_freeform(shape, transform, source)
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
            return self._render_text_or_shape(shape, transform, source)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return self._render_auto_shape(shape, transform, source)
        return ""

    def _render_group(self, shape, transform: Transform, source: str) -> str:
        xfrm = self._xfrm(shape)
        if xfrm is None:
            return ""
        off_x, off_y, ext_x, ext_y, ch_off_x, ch_off_y, ch_ext_x, ch_ext_y = xfrm
        if not ch_ext_x or not ch_ext_y:
            return ""

        child_sx = transform.sx * ext_x / ch_ext_x
        child_sy = transform.sy * ext_y / ch_ext_y
        child_transform = Transform(
            ox=transform.ox + off_x * transform.sx - ch_off_x * child_sx,
            oy=transform.oy + off_y * transform.sy - ch_off_y * child_sy,
            sx=child_sx,
            sy=child_sy,
        )
        return "\n".join(
            rendered
            for rendered in (self._render_shape(child, child_transform, source) for child in shape.shapes)
            if rendered
        )

    def _render_picture(self, shape, transform: Transform, source: str) -> str:
        rect = self._shape_rect_px(shape, transform)
        if rect is None:
            return ""
        asset = self._picture_asset(shape)
        if asset is None:
            return ""
        style = self._box_style(rect, shape)
        crop = self._crop(shape)
        name = html.escape(getattr(shape, "name", "picture"))
        if crop:
            l, t, r, b = crop
            visible_w = max(0.01, 1 - l - r)
            visible_h = max(0.01, 1 - t - b)
            img_style = (
                f"left:{-l / visible_w * 100:.4f}%;top:{-t / visible_h * 100:.4f}%;"
                f"width:{100 / visible_w:.4f}%;height:{100 / visible_h:.4f}%;"
            )
        else:
            img_style = "left:0;top:0;width:100%;height:100%;"
        return (
            f'<div class="ls-element ls-picture" data-source="{source}" style="{style}">'
            f'<img data-src="{html.escape(asset["relative"])}" alt="{name}" loading="lazy" '
            f'decoding="async" style="{img_style}"></div>'
        )

    def _render_text_or_shape(self, shape, transform: Transform, source: str) -> str:
        rect = self._shape_rect_px(shape, transform)
        if rect is None:
            return ""
        shape_style = self._shape_style(shape)
        text_html = self._text_html(shape)
        style = self._box_style(rect, shape) + shape_style + self._text_box_style(shape)
        cls = "ls-shape ls-text" if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE else "ls-text"
        return f'<div class="ls-element {cls}" data-source="{source}" style="{style}">{text_html}</div>'

    def _render_auto_shape(self, shape, transform: Transform, source: str) -> str:
        rect = self._shape_rect_px(shape, transform)
        if rect is None:
            return ""
        style = self._box_style(rect, shape) + self._shape_style(shape)
        preset = self._preset_geometry(shape)
        if preset == "ellipse":
            style += "border-radius:50%;"
        elif preset == "roundRect":
            style += "border-radius:24px;"
        return f'<div class="ls-element ls-shape" data-source="{source}" style="{style}"></div>'

    def _render_line(self, shape, transform: Transform, source: str) -> str:
        rect = self._shape_rect_px(shape, transform)
        if rect is None:
            return ""
        x, y, w, h = rect
        sp_pr = self._sp_pr(shape)
        stroke, width, alpha = self._line_style(sp_pr)
        if stroke == "none":
            return ""
        flip_h = self._has_xfrm_flag(shape, "flipH")
        flip_v = self._has_xfrm_flag(shape, "flipV")
        x1, y1 = (w, 0) if flip_h else (0, 0)
        x2, y2 = (0, h) if flip_h else (w, h)
        if flip_v:
            y1, y2 = h - y1, h - y2
        svg = (
            f'<svg viewBox="0 0 {max(w, 1):.3f} {max(h, 1):.3f}" preserveAspectRatio="none" '
            f'width="100%" height="100%"><line x1="{x1:.3f}" y1="{y1:.3f}" x2="{x2:.3f}" y2="{y2:.3f}" '
            f'stroke="{stroke}" stroke-width="{max(width, 1):.3f}" stroke-opacity="{alpha:.3f}" '
            f'stroke-linecap="round"/></svg>'
        )
        return (
            f'<div class="ls-element ls-vector" data-source="{source}" '
            f'style="left:{x:.3f}px;top:{y:.3f}px;width:{max(w, 1):.3f}px;height:{max(h, 1):.3f}px;">{svg}</div>'
        )

    def _render_freeform(self, shape, transform: Transform, source: str) -> str:
        rect = self._shape_rect_px(shape, transform)
        if rect is None:
            return ""
        paths = self._svg_paths(shape)
        if not paths:
            return ""
        sp_pr = self._sp_pr(shape)
        fill, fill_alpha = self._fill_style(sp_pr)
        stroke, stroke_width, stroke_alpha = self._line_style(sp_pr)
        x, y, w, h = rect
        body = []
        for view_w, view_h, path_d in paths:
            body.append(
                f'<svg viewBox="0 0 {view_w} {view_h}" preserveAspectRatio="none" width="100%" height="100%">'
                f'<path d="{html.escape(path_d)}" fill="{fill}" fill-opacity="{fill_alpha:.3f}" '
                f'stroke="{stroke}" stroke-width="{stroke_width:.3f}" stroke-opacity="{stroke_alpha:.3f}"/></svg>'
            )
        return (
            f'<div class="ls-element ls-vector" data-source="{source}" '
            f'style="left:{x:.3f}px;top:{y:.3f}px;width:{w:.3f}px;height:{h:.3f}px;">'
            f'{"".join(body)}</div>'
        )

    def _text_html(self, shape) -> str:
        paragraphs = []
        for para in shape.text_frame.paragraphs:
            runs = []
            for run in para.runs:
                if run.text == "":
                    continue
                runs.append(f'<span style="{self._run_style(run, para, shape)}">{html.escape(run.text)}</span>')
            if not runs and para.text:
                runs.append(f"<span>{html.escape(para.text)}</span>")
            align = self._paragraph_align(para)
            margin = self._paragraph_margin(para)
            paragraphs.append(f'<p style="text-align:{align};{margin}">{"".join(runs)}</p>')
        return "".join(paragraphs)

    def _run_style(self, run, para, shape) -> str:
        font = run.font
        styles = []
        if font.name:
            styles.append(f"font-family:'{self._css_string(font.name)}', {FONT_FALLBACK};")
        else:
            styles.append(f"font-family:{FONT_FALLBACK};")
        if font.size:
            styles.append(f"font-size:{font.size.pt * self._font_px_ratio():.3f}px;")
        color, alpha = self._run_color(run, para, shape)
        styles.append(f"color:{color};")
        if alpha < 1:
            styles.append(f"opacity:{alpha:.3f};")
        if font.bold:
            styles.append("font-weight:700;")
        if font.italic:
            styles.append("font-style:italic;")
        return "".join(styles)

    def _text_box_style(self, shape) -> str:
        body_pr = shape._element.find(f".//{TAG_A}bodyPr")
        if body_pr is None:
            return ""
        styles = []
        if body_pr.get("wrap") == "none":
            styles.append("white-space:pre;")
        if body_pr.get("anchor") == "ctr":
            styles.append("justify-content:center;")
        elif body_pr.get("anchor") == "b":
            styles.append("justify-content:flex-end;")
        return "".join(styles)

    def _paragraph_align(self, para) -> str:
        value = str(para.alignment or "").lower()
        if "center" in value:
            return "center"
        if "right" in value:
            return "right"
        if "justify" in value:
            return "justify"
        return "left"

    def _paragraph_margin(self, para) -> str:
        level = getattr(para, "level", 0) or 0
        if level <= 0:
            return ""
        return f"padding-left:{level * 28}px;"

    def _shape_style(self, shape) -> str:
        sp_pr = self._sp_pr(shape)
        fill, fill_alpha = self._fill_style(sp_pr)
        stroke, stroke_width, stroke_alpha = self._line_style(sp_pr)
        styles = []
        if fill != "none":
            styles.append(f"background:{fill};")
            if fill_alpha < 1:
                styles.append(f"background-color:{fill};")
                styles.append(f"opacity:{fill_alpha:.3f};")
        else:
            styles.append("background:transparent;")
        if stroke != "none" and stroke_width > 0:
            rgba = self._rgba(stroke, stroke_alpha)
            styles.append(f"border:{stroke_width:.3f}px solid {rgba};")
        return "".join(styles)

    def _box_style(self, rect: tuple[float, float, float, float], shape) -> str:
        x, y, w, h = rect
        styles = [f"left:{x:.3f}px;top:{y:.3f}px;width:{w:.3f}px;height:{h:.3f}px;"]
        rotation = getattr(shape, "rotation", 0) or 0
        if rotation:
            styles.append(f"transform:rotate({rotation:.3f}deg);transform-origin:center center;")
        return "".join(styles)

    def _shape_rect_px(self, shape, transform: Transform) -> tuple[float, float, float, float] | None:
        xfrm = self._xfrm(shape)
        if xfrm is None:
            return None
        x, y, w, h, *_ = xfrm
        tx, ty, tw, th = transform.apply_rect(x, y, w, h)
        return tx * self.scale_x, ty * self.scale_y, tw * self.scale_x, th * self.scale_y

    def _xfrm(self, shape) -> tuple[int, int, int, int, int, int, int, int] | None:
        sp_pr = self._sp_pr(shape)
        if sp_pr is None:
            return None
        xfrm = sp_pr.find(f"{TAG_A}xfrm")
        if xfrm is None:
            return None
        off = xfrm.find(f"{TAG_A}off")
        ext = xfrm.find(f"{TAG_A}ext")
        ch_off = xfrm.find(f"{TAG_A}chOff")
        ch_ext = xfrm.find(f"{TAG_A}chExt")
        if off is None or ext is None:
            return None
        return (
            int(off.get("x", "0")),
            int(off.get("y", "0")),
            int(ext.get("cx", "0")),
            int(ext.get("cy", "0")),
            int(ch_off.get("x", "0")) if ch_off is not None else 0,
            int(ch_off.get("y", "0")) if ch_off is not None else 0,
            int(ch_ext.get("cx", ext.get("cx", "0"))) if ch_ext is not None else int(ext.get("cx", "0")),
            int(ch_ext.get("cy", ext.get("cy", "0"))) if ch_ext is not None else int(ext.get("cy", "0")),
        )

    def _sp_pr(self, shape):
        elm = shape._element
        for tag in ("spPr", "grpSpPr", "picPr"):
            child = elm.find(f"{TAG_P}{tag}")
            if child is not None:
                return child
        return None

    def _fill_style(self, sp_pr) -> tuple[str, float]:
        if sp_pr is None:
            return "none", 1.0
        if sp_pr.find(f"{TAG_A}noFill") is not None:
            return "none", 1.0
        solid = sp_pr.find(f"{TAG_A}solidFill")
        if solid is None:
            return "none", 1.0
        return self._color_from_solid(solid)

    def _line_style(self, sp_pr) -> tuple[str, float, float]:
        if sp_pr is None:
            return "none", 0.0, 1.0
        ln = sp_pr.find(f"{TAG_A}ln")
        if ln is None or ln.find(f"{TAG_A}noFill") is not None:
            return "none", 0.0, 1.0
        width = int(ln.get("w", "9525")) * self.scale_x
        solid = ln.find(f"{TAG_A}solidFill")
        if solid is None:
            return "none", 0.0, 1.0
        color, alpha = self._color_from_solid(solid)
        return color, width, alpha

    def _run_color(self, run, para, shape) -> tuple[str, float]:
        r_pr = run._r.find(f"{TAG_A}rPr")
        if r_pr is not None:
            solid = r_pr.find(f"{TAG_A}solidFill")
            if solid is not None:
                return self._color_from_solid(solid)
        p_pr = para._p.find(f"{TAG_A}pPr")
        if p_pr is not None:
            solid = p_pr.find(f"{TAG_A}defRPr/{TAG_A}solidFill")
            if solid is not None:
                return self._color_from_solid(solid)
        level = getattr(para, "level", 0) or 0
        lvl_pr = shape._element.find(
            f".//{TAG_A}lstStyle/{TAG_A}lvl{level + 1}pPr/{TAG_A}defRPr/{TAG_A}solidFill"
        )
        if lvl_pr is not None:
            return self._color_from_solid(lvl_pr)
        return "#000000", 1.0

    def _color_from_solid(self, solid) -> tuple[str, float]:
        srgb = solid.find(f"{TAG_A}srgbClr")
        scheme = solid.find(f"{TAG_A}schemeClr")
        sys = solid.find(f"{TAG_A}sysClr")
        alpha = self._alpha(solid)
        if srgb is not None:
            return self._apply_modifiers("#" + srgb.get("val", "000000"), srgb), alpha * self._alpha(srgb)
        if sys is not None:
            return "#" + sys.get("lastClr", "000000"), alpha * self._alpha(sys)
        if scheme is not None:
            base = self.theme.get(scheme.get("val", ""), "#000000")
            return self._apply_modifiers(base, scheme), alpha * self._alpha(scheme)
        return "#000000", alpha

    def _apply_modifiers(self, color: str, elm) -> str:
        rgb = [int(color[i : i + 2], 16) for i in (1, 3, 5)]
        tint = elm.find(f"{TAG_A}tint")
        shade = elm.find(f"{TAG_A}shade")
        lum_mod = elm.find(f"{TAG_A}lumMod")
        lum_off = elm.find(f"{TAG_A}lumOff")
        if shade is not None:
            factor = int(shade.get("val", "100000")) / 100000
            rgb = [round(c * factor) for c in rgb]
        if tint is not None:
            factor = int(tint.get("val", "0")) / 100000
            rgb = [round(c + (255 - c) * factor) for c in rgb]
        if lum_mod is not None:
            factor = int(lum_mod.get("val", "100000")) / 100000
            rgb = [round(c * factor) for c in rgb]
        if lum_off is not None:
            offset = int(lum_off.get("val", "0")) / 100000
            rgb = [round(c + 255 * offset) for c in rgb]
        rgb = [max(0, min(255, c)) for c in rgb]
        return "#" + "".join(f"{c:02x}" for c in rgb)

    def _alpha(self, elm) -> float:
        alpha = elm.find(f"{TAG_A}alpha")
        if alpha is None:
            return 1.0
        return max(0.0, min(1.0, int(alpha.get("val", "100000")) / 100000))

    def _rgba(self, color: str, alpha: float) -> str:
        if alpha >= 1 or color == "none":
            return color
        rgb = tuple(int(color[i : i + 2], 16) for i in (1, 3, 5))
        return f"rgba({rgb[0]},{rgb[1]},{rgb[2]},{alpha:.3f})"

    def _preset_geometry(self, shape) -> str:
        sp_pr = self._sp_pr(shape)
        if sp_pr is None:
            return "rect"
        prst = sp_pr.find(f"{TAG_A}prstGeom")
        return prst.get("prst", "rect") if prst is not None else "rect"

    def _svg_paths(self, shape) -> list[tuple[int, int, str]]:
        sp_pr = self._sp_pr(shape)
        if sp_pr is None:
            return []
        paths = []
        for path in sp_pr.findall(f".//{TAG_A}path"):
            view_w = int(path.get("w", "1"))
            view_h = int(path.get("h", "1"))
            commands = []
            for child in list(path):
                tag = self._local_name(child.tag)
                points = child.findall(f"{TAG_A}pt")
                if tag == "moveTo" and points:
                    commands.append(f'M {points[0].get("x")} {points[0].get("y")}')
                elif tag == "lnTo" and points:
                    commands.append(f'L {points[0].get("x")} {points[0].get("y")}')
                elif tag == "cubicBezTo" and len(points) == 3:
                    commands.append(
                        "C "
                        + " ".join(f'{pt.get("x")} {pt.get("y")}' for pt in points)
                    )
                elif tag == "quadBezTo" and len(points) == 2:
                    commands.append(
                        "Q "
                        + " ".join(f'{pt.get("x")} {pt.get("y")}' for pt in points)
                    )
                elif tag == "close":
                    commands.append("Z")
            if commands:
                paths.append((view_w, view_h, " ".join(commands)))
        return paths

    def _picture_asset(self, shape) -> dict | None:
        blip = shape._element.find(f".//{TAG_A}blip")
        if blip is None:
            return None
        rid = blip.get(f"{TAG_R}embed")
        if not rid:
            return None
        image_part = shape.part.related_part(rid)
        blob = image_part.blob
        digest = hashlib.sha1(blob).hexdigest()[:16]
        ext = Path(image_part.partname).suffix.lower() or ".bin"
        converted = self._browser_image(blob, ext, digest)
        if converted is not None:
            blob, ext = converted
        key = digest + ext
        if key not in self.asset_by_hash:
            name = f"{digest}{ext}"
            target = self.assets_dir / name
            target.write_bytes(blob)
            entry = {
                "path": str(target.relative_to(self.out_dir.parent)),
                "relative": f"assets/pptx-media/{name}",
                "bytes": len(blob),
            }
            self.asset_by_hash[key] = entry
            self.assets.append({"path": entry["path"], "bytes": len(blob)})
        return self.asset_by_hash[key]

    def _browser_image(self, blob: bytes, ext: str, digest: str) -> tuple[bytes, str] | None:
        if ext not in {".emf", ".wmf"}:
            return None
        png_blob = self._metafile_to_png(blob, digest)
        if png_blob:
            return png_blob, ".png"
        thumb_blob = self._xmp_thumbnail(blob)
        if thumb_blob:
            return thumb_blob, ".jpg"
        return None

    def _metafile_to_png(self, blob: bytes, digest: str) -> bytes | None:
        pdf_offset = blob.find(b"%PDF-")
        if pdf_offset < 0:
            return None
        with tempfile.TemporaryDirectory(prefix=f"lark-slides-{digest}-") as tmp:
            tmp_dir = Path(tmp)
            pdf_path = tmp_dir / "source.pdf"
            png_path = tmp_dir / "source.png"
            pdf_path.write_bytes(blob[pdf_offset:])
            result = subprocess.run(
                ["sips", "-s", "format", "png", str(pdf_path), "--out", str(png_path)],
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                check=False,
            )
            if result.returncode == 0 and png_path.exists():
                return png_path.read_bytes()
        return None

    def _xmp_thumbnail(self, blob: bytes) -> bytes | None:
        match = re.search(rb"<xmpGImg:image>(.*?)</xmpGImg:image>", blob, re.S)
        if not match:
            return None
        encoded = match.group(1).replace(b"&#xA;", b"").replace(b"\n", b"").replace(b"\r", b"")
        try:
            return base64.b64decode(encoded)
        except Exception:
            return None

    def _crop(self, shape) -> tuple[float, float, float, float] | None:
        src_rect = shape._element.find(f".//{TAG_A}srcRect")
        if src_rect is None:
            return None
        values = []
        for key in ("l", "t", "r", "b"):
            values.append(int(src_rect.get(key, "0")) / 100000)
        if not any(values):
            return None
        return tuple(values)

    def _has_xfrm_flag(self, shape, flag: str) -> bool:
        sp_pr = self._sp_pr(shape)
        if sp_pr is None:
            return False
        xfrm = sp_pr.find(f"{TAG_A}xfrm")
        return xfrm is not None and xfrm.get(flag) == "1"

    def _is_placeholder(self, shape) -> bool:
        try:
            return shape.is_placeholder
        except Exception:
            return False

    def _is_hidden(self, shape) -> bool:
        c_nv_pr = shape._element.find(f".//{TAG_P}cNvPr")
        return c_nv_pr is not None and c_nv_pr.get("hidden") == "1"

    def _slide_title(self, slide, index: int) -> str:
        texts = []

        def collect(shapes):
            for shape in shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    collect(shape.shapes)
                elif getattr(shape, "has_text_frame", False):
                    text = re.sub(r"\s+", " ", shape.text_frame.text).strip()
                    if text:
                        texts.append(text)

        collect(slide.shapes)
        return texts[0][:80] if texts else f"Slide {index}"

    def _slide_section(self, index: int, title: str, body: str) -> str:
        safe_title = html.escape(title, quote=True)
        return (
            f'<section class="ls-slide" data-slide="{index}" data-title="{safe_title}" '
            f'aria-label="{safe_title}">\n'
            f'<div class="ls-slide-inner">\n{body}\n</div>\n'
            f"</section>"
        )

    def _deck_html(self, slides: str) -> str:
        return f"""<!doctype html>
<html lang="zh-CN">
  <head>
    <meta charset="utf-8" />
    <meta name="viewport" content="width=device-width, initial-scale=1" />
    <title>Lark Design Guidelines HTML PPT</title>
    <link rel="icon" href="data:," />
    <link rel="stylesheet" href="../sdk/lark-slides.css" />
  </head>
  <body>
    <div class="ls-app" data-lark-deck data-title="Lark Design Guidelines">
      <div class="ls-stage">
{slides}
      </div>
    </div>
    <script src="../sdk/lark-slides.js"></script>
    <script>
      window.LarkSlides.createDeck({{ mount: document.querySelector("[data-lark-deck]") }});
    </script>
  </body>
</html>
"""

    def _font_px_ratio(self) -> float:
        slide_width_points = self.prs.slide_width / 12700
        return CANVAS["width"] / slide_width_points

    def _load_theme_colors(self) -> dict[str, str]:
        colors = {
            "dk1": "#000000",
            "lt1": "#ffffff",
            "dk2": "#1f2329",
            "lt2": "#f5f7fa",
            "tx1": "#000000",
            "bg1": "#ffffff",
            "tx2": "#1f2329",
            "bg2": "#f5f7fa",
        }
        with zipfile.ZipFile(self.pptx_path) as archive:
            theme_names = [name for name in archive.namelist() if name.startswith("ppt/theme/theme") and name.endswith(".xml")]
            if not theme_names:
                return colors
            root = ET.fromstring(archive.read(theme_names[0]))
        scheme = root.find(".//a:clrScheme", NS)
        if scheme is None:
            return colors
        for child in list(scheme):
            name = self._local_name(child.tag)
            srgb = child.find("a:srgbClr", NS)
            sys = child.find("a:sysClr", NS)
            if srgb is not None:
                colors[name] = "#" + srgb.attrib.get("val", "000000")
            elif sys is not None:
                colors[name] = "#" + sys.attrib.get("lastClr", "000000")
        colors["tx1"] = colors.get("dk1", colors["tx1"])
        colors["bg1"] = colors.get("lt1", colors["bg1"])
        colors["tx2"] = colors.get("dk2", colors["tx2"])
        colors["bg2"] = colors.get("lt2", colors["bg2"])
        return colors

    def _local_name(self, tag: str) -> str:
        return tag.split("}", 1)[-1]

    def _css_string(self, value: str) -> str:
        return value.replace("\\", "\\\\").replace("'", "\\'")


def main() -> None:
    parser = argparse.ArgumentParser(description="Convert a PPTX deck into an HTML slide deck.")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--out", type=Path, default=Path("dist"))
    args = parser.parse_args()

    Converter(args.pptx.resolve(), args.out.resolve()).run()


if __name__ == "__main__":
    main()
